import os
import re
import json
import numpy as np
from pypdf import PdfReader
from google import genai
from google.genai import types
from openai import OpenAI

class DocumentExtractor:
    """Extracts text page by page or slide by slide from various document types (.pdf, .docx, .pptx, .txt, .md)."""
    @staticmethod
    def extract_text_with_pages(file_path: str) -> list[dict]:
        ext = os.path.splitext(file_path)[1].lower()
        pages = []
        
        if ext == ".pdf":
            try:
                reader = PdfReader(file_path)
                for i, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    # Clean extra whitespace
                    text = re.sub(r'\s+', ' ', text).strip()
                    if text:
                        pages.append({
                            "page_number": i + 1,
                            "text": text
                        })
            except Exception as e:
                print(f"Error reading PDF {file_path}: {e}")
                raise e
                
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            cleaned_text = re.sub(r'\s+', ' ', shape.text).strip()
                            if cleaned_text:
                                slide_text.append(cleaned_text)
                    full_text = "\n".join(slide_text)
                    if full_text.strip():
                        pages.append({
                            "page_number": i + 1,
                            "text": full_text
                        })
            except Exception as e:
                print(f"Error reading PowerPoint PPTX {file_path}: {e}")
                raise e
                
        elif ext == ".docx":
            try:
                import docx
                doc = docx.Document(file_path)
                paragraphs = []
                for p in doc.paragraphs:
                    if p.text.strip():
                        paragraphs.append(p.text.strip())
                full_text = "\n".join(paragraphs)
                
                # Assign virtual pages every ~3000 characters
                page_size = 3000
                total_len = len(full_text)
                if total_len > 0:
                    num_pages = (total_len + page_size - 1) // page_size
                    for i in range(num_pages):
                        start = i * page_size
                        end = min(start + page_size, total_len)
                        pages.append({
                            "page_number": i + 1,
                            "text": full_text[start:end]
                        })
            except Exception as e:
                print(f"Error reading Word DOCX {file_path}: {e}")
                raise e
                
        elif ext in [".txt", ".md"]:
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    full_text = f.read()
                full_text = re.sub(r'\s+', ' ', full_text).strip()
                
                # Assign virtual pages every ~3000 characters
                page_size = 3000
                total_len = len(full_text)
                if total_len > 0:
                    num_pages = (total_len + page_size - 1) // page_size
                    for i in range(num_pages):
                        start = i * page_size
                        end = min(start + page_size, total_len)
                        pages.append({
                            "page_number": i + 1,
                            "text": full_text[start:end]
                        })
            except Exception as e:
                print(f"Error reading text file {file_path}: {e}")
                raise e
        else:
            raise ValueError(f"Unsupported file format: {ext}")
            
        return pages


class TextChunker:
    """Splits document pages into smaller semantic chunks with overlapping boundaries."""
    @staticmethod
    def split_pages(pages: list[dict], chunk_size: int = 800, chunk_overlap: int = 150) -> list[dict]:
        chunks = []
        chunk_idx = 0
        
        for page in pages:
            page_num = page["page_number"]
            text = page["text"]
            
            # Use a robust sliding window split by words
            words = text.split(" ")
            current_chunk_words = []
            current_len = 0
            
            for word in words:
                word_len = len(word) + 1
                if current_len + word_len > chunk_size and current_chunk_words:
                    chunk_text = " ".join(current_chunk_words).strip()
                    if chunk_text:
                        chunks.append({
                            "chunk_id": f"c_{chunk_idx}",
                            "text": chunk_text,
                            "page": page_num
                        })
                        chunk_idx += 1
                    
                    # Compute word overlap
                    overlap_words = []
                    overlap_len = 0
                    for w in reversed(current_chunk_words):
                        if overlap_len + len(w) + 1 > chunk_overlap:
                            break
                        overlap_words.insert(0, w)
                        overlap_len += len(w) + 1
                    
                    current_chunk_words = overlap_words
                    current_len = overlap_len
                
                current_chunk_words.append(word)
                current_len += word_len
                
            # Grab any remaining text on the page
            if current_chunk_words:
                chunk_text = " ".join(current_chunk_words).strip()
                if chunk_text:
                    chunks.append({
                        "chunk_id": f"c_{chunk_idx}",
                        "text": chunk_text,
                        "page": page_num
                    })
                    chunk_idx += 1
                    
        return chunks

class NumPyVectorStore:
    """A pure NumPy in-memory vector store for cosine similarity matching."""
    def __init__(self, storage_dir: str = "storage"):
        self.storage_dir = storage_dir
        os.makedirs(storage_dir, exist_ok=True)
        self.db_file = os.path.join(storage_dir, "vector_store.json")
        self.chunks = []      # list of dict: {chunk_id, text, page, doc_id}
        self.embeddings = []  # list of np.ndarray
        self.load()

    def add_document(self, doc_id: str, chunks: list[dict], embeddings: list[list[float]]):
        # Remove any existing chunks for this doc_id to support re-upload
        self.remove_document(doc_id)
        
        for chunk, emb in zip(chunks, embeddings):
            chunk_data = {**chunk, "doc_id": doc_id}
            self.chunks.append(chunk_data)
            self.embeddings.append(np.array(emb, dtype=np.float32))
        self.save()

    def remove_document(self, doc_id: str):
        indices_to_keep = [i for i, chunk in enumerate(self.chunks) if chunk["doc_id"] != doc_id]
        self.chunks = [self.chunks[i] for i in indices_to_keep]
        self.embeddings = [self.embeddings[i] for i in indices_to_keep]
        self.save()

    def search(self, query_emb: list[float], doc_id: str, top_k: int = 4) -> list[dict]:
        if not self.embeddings:
            return []
            
        filtered_indices = [i for i, chunk in enumerate(self.chunks) if chunk["doc_id"] == doc_id]
        if not filtered_indices:
            return []
            
        q_vec = np.array(query_emb, dtype=np.float32)
        q_norm = np.linalg.norm(q_vec)
        if q_norm == 0:
            return []
            
        results = []
        for idx in filtered_indices:
            emb = self.embeddings[idx]
            emb_norm = np.linalg.norm(emb)
            if emb_norm == 0:
                score = 0.0
            else:
                score = float(np.dot(q_vec, emb) / (q_norm * emb_norm))
            results.append((score, self.chunks[idx]))
            
        results.sort(key=lambda x: x[0], reverse=True)
        return [{"score": score, **chunk} for score, chunk in results[:top_k]]

    def save(self):
        """Serializes the index to local JSON for persistence between server reloads."""
        serialized_embeddings = [emb.tolist() for emb in self.embeddings]
        data = {
            "chunks": self.chunks,
            "embeddings": serialized_embeddings
        }
        with open(self.db_file, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

    def load(self):
        """Loads index from JSON on startup."""
        if os.path.exists(self.db_file):
            try:
                with open(self.db_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                self.chunks = data.get("chunks", [])
                self.embeddings = [np.array(emb, dtype=np.float32) for emb in data.get("embeddings", [])]
            except Exception as e:
                print(f"Failed to load vector store: {e}. Starting fresh.")
                self.chunks = []
                self.embeddings = []

class LLMService:
    """Handles communication with Google Gemini and OpenAI for embeddings and response generation."""
    @staticmethod
    def get_embedding(text: str, provider: str, api_key: str) -> list[float]:
        if provider == "openai":
            client = OpenAI(api_key=api_key)
            response = client.embeddings.create(
                model="text-embedding-3-small",
                input=[text]
            )
            return response.data[0].embedding
        else: # Default to Gemini
            # Use google-genai SDK
            client = genai.Client(api_key=api_key)
            response = client.models.embed_content(
                model="gemini-embedding-004",
                contents=text
            )
            # Handle list vs single response in embeddings API
            if isinstance(response.embeddings, list):
                return response.embeddings[0].values
            return response.embeddings.values

    @staticmethod
    def get_batch_embeddings(texts: list[str], provider: str, api_key: str) -> list[list[float]]:
        embeddings = []
        if provider == "openai":
            client = OpenAI(api_key=api_key)
            # Batch size limits in OpenAI are generous, but we split just in case
            batch_size = 100
            for i in range(0, len(texts), batch_size):
                batch = texts[i:i+batch_size]
                response = client.embeddings.create(
                    model="text-embedding-3-small",
                    input=batch
                )
                embeddings.extend([item.embedding for item in response.data])
        else:
            client = genai.Client(api_key=api_key)
            # google-genai supports batch embedding via list contents
            # We will embed in chunks of 50 to prevent size limits
            batch_size = 50
            for i in range(0, len(texts), batch_size):
                batch = texts[i:i+batch_size]
                response = client.models.embed_content(
                    model="gemini-embedding-004",
                    contents=batch
                )
                for emb in response.embeddings:
                    embeddings.append(emb.values)
        return embeddings

    @staticmethod
    def generate_answer(query: str, context: str, provider: str, api_key: str) -> str:
        system_prompt = (
            "You are an expert academic AI Study Assistant. Answer the student's question based strictly "
            "on the provided context extracted from their uploaded lecture PDF notes. If the answer cannot "
            "be found or inferred from the context, state clearly that the notes do not contain this information. "
            "Do not make up facts outside the provided context. Include helpful references to the Page number "
            "when referring to specific points in your response.\n\n"
            f"--- CONTEXT START ---\n{context}\n--- CONTEXT END ---\n"
        )
        
        if provider == "openai":
            client = OpenAI(api_key=api_key)
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": query}
                ],
                temperature=0.3
            )
            return response.choices[0].message.content
        else:
            client = genai.Client(api_key=api_key)
            response = client.models.generate_content(
                model="gemini-2.5-flash", # or gemini-2.5-flash
                contents=query,
                config=types.GenerateContentConfig(
                    system_instruction=system_prompt,
                    temperature=0.3
                )
            )
            return response.text
